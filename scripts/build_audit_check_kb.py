from __future__ import annotations

import re
import zipfile
from collections import defaultdict
from pathlib import Path
from xml.etree import ElementTree as ET

import openpyxl
import pdfplumber
from docx import Document
from pptx import Presentation


ROOT = Path(r"Y:\IFC_V\公用區\115年感管查核")
OUT = Path(__file__).resolve().parents[1] / "output" / "coze_upload" / "115年感管查核_臨床同仁重點.md"

SKIP_TERMS = (
    "座位表",
    "簽到",
    "路線",
    "程序表",
    "便簽",
    "來文",
    "分工表",
    "負責人名單",
    "回條",
    "Thumbs",
)

PREFERRED_TERMS = (
    "全院宣導",
    "護理部宣導",
    "查核基準",
    "評分說明",
    "作業手冊",
    "自評表",
    "追蹤訪查",
    "查核結果",
    "意見表",
    "健康監測",
    "供應室",
    "HAI",
    "海報",
)

TOPICS: dict[str, list[str]] = {
    "手部衛生": ["手部衛生", "洗手", "乾洗手", "濕洗手", "酒精性乾洗手", "五時機", "手套"],
    "個人防護裝備與隔離": ["隔離", "接觸", "飛沫", "空氣", "N95", "口罩", "面罩", "護目鏡", "防護衣", "PPE", "負壓"],
    "侵入性裝置感染預防": ["導尿管", "尿管", "中心導管", "CVC", "CLABSI", "CAUTI", "呼吸器", "VAP", "管路", "留置"],
    "抗生素管理": ["抗生素", "抗菌藥物", "抗藥", "DDD", "管制性抗生素", "預防性抗生素", "AST"],
    "MDRO與抗藥菌": ["MDRO", "MRSA", "VRE", "CRE", "CRAB", "CRPA", "抗藥性", "多重抗藥", "菌株"],
    "環境清潔與消毒": ["環境清潔", "環境消毒", "清消", "漂白水", "消毒劑", "高頻接觸", "終期清潔", "病室"],
    "醫材清洗消毒滅菌與供應室": ["消毒", "滅菌", "高層次消毒", "內視鏡", "器械", "衛材", "供應室", "包裝", "滅菌鍋"],
    "注射與尖銳物安全": ["針扎", "尖銳物", "血液體液", "暴露", "注射", "採血", "安全針具"],
    "員工健康與健康監測": ["健康監測", "員工", "症狀通報", "疫苗", "流感疫苗", "B型肝炎", "胸部X光", "職業暴露"],
    "法定傳染病與疫情通報": ["法定傳染病", "通報", "TOCC", "疫區", "旅遊史", "群聚", "疑似", "院內群聚"],
    "醫療照護相關感染監測": ["醫療照護相關感染", "HAI", "感染密度", "監測", "THAS", "感染率", "群突發"],
    "病人安置與動線": ["病人安置", "床位", "單人房", "轉送", "檢查", "動線", "候診", "急診"],
    "廢棄物與布服": ["廢棄物", "感染性廢棄物", "布服", "污衣", "垃圾", "污物", "尖銳物收集盒"],
    "文件與現場查核回答": ["查核", "自評", "佐證", "紀錄", "稽核", "改善", "教育訓練", "SOP", "政策"],
}


def normalize(text: str) -> str:
    text = text.replace("\u3000", " ").replace("\xa0", " ")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def extract_docx(path: Path) -> str:
    try:
        doc = Document(str(path))
        parts: list[str] = []
        for p in doc.paragraphs:
            if p.text.strip():
                parts.append(p.text.strip())
        for table in doc.tables:
            for row in table.rows:
                row_text = "；".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    parts.append(row_text)
        return "\n".join(parts)
    except Exception:
        return ""


def extract_doc_legacy(path: Path) -> str:
    # Legacy .doc files are skipped unless Word-converted copies exist elsewhere.
    return ""


def extract_pptx(path: Path) -> str:
    try:
        prs = Presentation(str(path))
        parts: list[str] = []
        for i, slide in enumerate(prs.slides, start=1):
            slide_parts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_parts.append(shape.text.strip())
            if slide_parts:
                parts.append(f"投影片{i}: " + "\n".join(slide_parts))
        return "\n".join(parts)
    except Exception:
        return ""


def extract_pdf(path: Path, max_pages: int = 80) -> str:
    try:
        parts: list[str] = []
        with pdfplumber.open(str(path)) as pdf:
            for page in pdf.pages[:max_pages]:
                text = page.extract_text() or ""
                if text.strip():
                    parts.append(text.strip())
        return "\n".join(parts)
    except Exception:
        return ""


def extract_xlsx(path: Path) -> str:
    try:
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        parts: list[str] = []
        for ws in wb.worksheets[:8]:
            rows = 0
            for row in ws.iter_rows(values_only=True):
                vals = [str(v).strip() for v in row if v is not None and str(v).strip()]
                if vals:
                    parts.append("；".join(vals))
                    rows += 1
                if rows >= 250:
                    break
        return "\n".join(parts)
    except Exception:
        return ""


def extract_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".docx":
        return extract_docx(path)
    if suffix == ".doc":
        return extract_doc_legacy(path)
    if suffix == ".pptx":
        return extract_pptx(path)
    if suffix == ".pdf":
        return extract_pdf(path)
    if suffix == ".xlsx":
        return extract_xlsx(path)
    return ""


def should_use(path: Path) -> bool:
    name = path.name
    if path.suffix.lower() not in {".pdf", ".docx", ".pptx", ".xlsx"}:
        return False
    if any(term in name for term in SKIP_TERMS):
        return False
    return any(term in str(path) for term in PREFERRED_TERMS)


def useful_lines(text: str) -> list[str]:
    lines: list[str] = []
    for raw in normalize(text).splitlines():
        line = raw.strip(" -•\t")
        if len(line) < 6 or len(line) > 180:
            continue
        if re.fullmatch(r"[\d\s./:()（）-]+", line):
            continue
        if re.search(r"(http|@|電話|分機|主席|座位|簽到|頁碼|第\s*\d+\s*頁)", line, re.I):
            continue
        if line.count(" ") > 25:
            continue
        lines.append(line)
    return lines


def score_line(line: str, terms: list[str]) -> int:
    score = sum(3 for term in terms if term.lower() in line.lower())
    score += sum(1 for word in ("應", "須", "需", "不可", "避免", "落實", "注意", "確認", "記錄", "通報", "隔離", "清潔") if word in line)
    return score


def collect() -> tuple[dict[str, list[tuple[int, str, str]]], list[Path]]:
    topic_hits: dict[str, list[tuple[int, str, str]]] = defaultdict(list)
    used_files: list[Path] = []
    for path in sorted(ROOT.rglob("*")):
        if not path.is_file() or not should_use(path):
            continue
        text = extract_text(path)
        if not text:
            continue
        used_files.append(path)
        lines = useful_lines(text)
        for topic, terms in TOPICS.items():
            for line in lines:
                score = score_line(line, terms)
                if score >= 4:
                    topic_hits[topic].append((score, line, str(path)))
    return topic_hits, used_files


def dedupe_hits(hits: list[tuple[int, str, str]], limit: int = 12) -> list[tuple[int, str, str]]:
    seen: set[str] = set()
    selected: list[tuple[int, str, str]] = []
    for score, line, source in sorted(hits, key=lambda x: (-x[0], len(x[1]))):
        key = re.sub(r"\d+", "", line)
        key = key[:60]
        if key in seen:
            continue
        seen.add(key)
        selected.append((score, line, source))
        if len(selected) >= limit:
            break
    return selected


def manual_guidance(topic: str) -> list[str]:
    guidance = {
        "手部衛生": [
            "進出病室、接觸病人前後、接觸病人體液或周邊環境後，都應落實手部衛生；戴手套不能取代手部衛生。",
            "查核時不只看設備是否存在，也會看臨床人員是否能在正確時機執行，並能說明乾洗手與濕洗手適用情境。",
        ],
        "個人防護裝備與隔離": [
            "隔離標示的目的不是貼疾病名稱，而是讓進入照護者知道應採取接觸、飛沫、空氣或特殊防護。",
            "進入隔離區前先確認 PPE 類型與穿脫順序；離開前應避免污染自己、環境與公共動線。",
        ],
        "侵入性裝置感染預防": [
            "導尿管、中心導管、呼吸器等裝置應每日評估必要性，能移除就移除，並維持照護 bundle 與紀錄完整。",
            "查核常看實際照護是否與紀錄一致，例如管路固定、標示、照護紀錄、感染徵象與移除評估。",
        ],
        "抗生素管理": [
            "抗生素使用應有適應症、培養檢體、去升階或停藥評估；不要把預防性或經驗性用藥當成無限期治療。",
            "被問到抗生素管理時，臨床端應能說明院內抗菌藥物管理、管制性抗生素申請與感染科/藥師介入機制。",
        ],
        "MDRO與抗藥菌": [
            "MDRO 病人重點是接觸隔離、手部衛生、器材專用或清消、檢查前通知與終期清潔，不是用醒目字樣標籤病人。",
            "再入院或既往陽性個案應依院內解隔與再篩檢流程處理；不能只憑病人主訴或單次陰性就解除。",
        ],
        "環境清潔與消毒": [
            "高頻接觸表面、病室設備、共用器材與終期清潔是查核重點；看的是現場是否做得到、紀錄是否能追溯。",
            "清消前先清潔可見髒污，再依疾病、污染情境與院內規範使用合適消毒劑與接觸時間。",
        ],
        "醫材清洗消毒滅菌與供應室": [
            "使用後器械應依污染程度、材質與用途進行清洗、消毒或滅菌；清潔不確實會影響後續消毒滅菌效果。",
            "滅菌包裝、有效期限、化學/生物監測、儲存環境與運送流程都屬臨床端應理解的病人安全環節。",
        ],
        "注射與尖銳物安全": [
            "抽血、注射、處置與尖銳物丟棄要遵守標準防護；針扎或血液體液暴露應立即處理並依院內流程通報。",
            "尖銳物收集盒不可過滿，使用後針具不可回套，採檢過程應避免造成工作人員與環境暴露。",
        ],
        "員工健康與健康監測": [
            "員工出現發燒、呼吸道、腸胃道、皮膚病灶或疑似傳染病暴露時，應依院內健康監測與通報流程處理。",
            "疫苗、胸部 X 光、暴露後追蹤與症狀通報不是行政作業而已，目的是避免院內傳播並保護病人與同仁。",
        ],
        "法定傳染病與疫情通報": [
            "遇疑似法定傳染病、群聚事件或具 TOCC 風險個案，臨床端應及早通報並確認檢體、隔離與病人動線。",
            "通報不是等確診才做；符合病例定義或疑似條件時應依時限與院內流程辦理。",
        ],
        "醫療照護相關感染監測": [
            "HAI 監測資料不只是感管中心報表，臨床單位應能理解感染率、裝置使用、抗藥菌與改善措施的關係。",
            "查核時可能追問單位是否知道自己的感染風險、是否有改善策略，以及改善後如何追蹤成效。",
        ],
        "病人安置與動線": [
            "有傳播風險的病人應依傳播途徑安排床位、檢查動線與候診方式；轉送或檢查前要通知接收單位。",
            "床位不足時應依風險分層與院內流程處理，不能只靠口頭提醒或把病人留在不適合的共同空間。",
        ],
        "廢棄物與布服": [
            "感染性廢棄物、尖銳物、污染布服與高風險檢體外包裝應依院內分類、包裝、標示與運送規定處理。",
            "布服或廢棄物處理重點是避免滲漏、飛散、刺傷與污染公共動線。",
        ],
        "文件與現場查核回答": [
            "查核回答應以現場實作為主，能說明自己單位平常怎麼做、在哪裡查 SOP、異常時通知誰。",
            "不要只背條文；若不確定，應回答會先依院內政策規章、單位主管與感染管制中心流程確認後執行。",
        ],
    }
    return guidance.get(topic, [])


def audit_prompts(topic: str) -> list[str]:
    prompts = {
        "手部衛生": [
            "查核常看：照護前後、無菌操作前、暴露體液風險後、接觸病人周邊環境後是否確實執行。",
            "同仁應能回答：什麼時候用乾洗手、什麼時候需要濕洗手；手部有明顯髒污時不可只用乾洗手。",
            "現場應注意：乾洗手液、洗手設備、擦手紙與補充狀態；工作車與隔離病室外是否方便取得手部衛生用品。",
        ],
        "個人防護裝備與隔離": [
            "查核常看：隔離標示是否清楚、PPE 是否取得容易、工作人員是否知道進出病室的穿脫順序。",
            "同仁應能回答：接觸隔離、飛沫隔離、空氣隔離各自的口罩、手套、隔離衣與病室要求。",
            "結核病或需空氣隔離個案，重點是負壓病室、N95 口罩、門關閉與解除隔離條件依院內規範確認。",
        ],
        "侵入性裝置感染預防": [
            "查核常看：導尿管、中心導管、呼吸器等裝置是否有每日必要性評估與照護紀錄。",
            "同仁應能回答：管路留置理由、移除評估、照護 bundle、感染徵象與異常通報方式。",
            "現場應注意：管路固定、引流袋位置、接頭清潔、敷料完整性與日期標示。",
        ],
        "抗生素管理": [
            "查核常看：手術預防性抗生素給藥時間、使用期間與超過時限時的病歷理由。",
            "同仁應能回答：二線以上或管制性抗生素需依院內機制申請或由主治醫師確認。",
            "臨床端應配合：培養採檢、抗生素 de-escalation、停藥評估與抗藥性趨勢回饋。",
        ],
        "MDRO與抗藥菌": [
            "查核常看：微生物報告是否能提示抗藥菌、病房是否依接觸隔離與環境清消落實。",
            "同仁應能回答：VRE、CRE、CRAB、CRPA、MRSA 等常見 MDRO 的隔離、檢查通知與終期清潔原則。",
            "觀念提醒：警示是為了讓照護者採取正確防護，不是把病人永久貼上危險標籤。",
        ],
        "環境清潔與消毒": [
            "查核常看：高頻接觸表面、病床周邊、護理站、遊戲室、廁所與隔離病室是否有清潔紀錄與稽核。",
            "同仁應能回答：先清潔可見髒污，再依疾病與污染情境消毒；漂白水濃度與接觸時間需依院內規範。",
            "現場應注意：共用器材用後清消、隔離病人出院或轉床後終期清潔、清潔用具分區使用。",
        ],
        "醫材清洗消毒滅菌與供應室": [
            "查核常看：器械是否先清洗再消毒/滅菌，滅菌包是否有完整標示、有效期限與監測紀錄。",
            "同仁應能回答：哪些物品需滅菌、哪些需高層次消毒、使用後器械如何送回供應室。",
            "現場應注意：清潔物與污染物分流，消毒或滅菌後物品應妥善包裝、儲存與運送。",
        ],
        "注射與尖銳物安全": [
            "查核常看：採血、注射、尖銳物丟棄是否符合標準防護，安全針具是否正確使用。",
            "同仁應能回答：血液體液暴露或針扎後，應立即沖洗、通報、就醫評估與追蹤。",
            "現場應注意：尖銳物收集盒不可過滿，不回套針頭，不將尖銳物留在工作車、床旁或垃圾袋內。",
        ],
        "員工健康與健康監測": [
            "查核常看：員工疫苗、胸部 X 光、健康監測、症狀通報與職業暴露追蹤是否落實。",
            "同仁應能回答：自己有發燒、呼吸道症狀、腸胃道症狀、皮膚病灶或疑似暴露時應如何通報。",
            "現場應注意：實習學生、志工與外包人員也可能被納入疫苗或健康監測提醒範圍。",
        ],
        "法定傳染病與疫情通報": [
            "查核常看：門急住診是否有 TOCC 詢問與紀錄，疑似法定傳染病是否依時限通報。",
            "同仁應能回答：發現疑似個案時，先隔離與通知，再確認通報、採檢、檢查動線與環境清消。",
            "群聚或新興傳染病事件時，應依院內應變計畫與感染管制中心指示啟動處理。",
        ],
        "醫療照護相關感染監測": [
            "查核常看：單位是否知道自己的 HAI、裝置相關感染、MDRO 與改善措施。",
            "同仁應能回答：感染率或監測數據不是只給感管中心看，臨床單位需依資料改善照護流程。",
            "現場應注意：若有異常趨勢或群突發，需有通報、調查、改善與追蹤紀錄。",
        ],
        "病人安置與動線": [
            "查核常看：急診、門診、病房、檢查單位是否能依傳播風險分流病人。",
            "同仁應能回答：隔離病人外出檢查前要通知接收單位，病人需戴口罩或採相應防護，轉送後器材與動線要清消。",
            "床位不足時，不自行降低防護等級；應依風險分層、主管與感染管制流程協調。",
        ],
        "廢棄物與布服": [
            "查核常看：感染性廢棄物、尖銳物、污染布服是否正確分類、包裝、加蓋與運送。",
            "同仁應能回答：污染布服或廢棄物不得任意堆放，應避免滲漏、飛散、刺傷與污染公共動線。",
            "現場應注意：感染性垃圾桶與污物桶應有合適容器與標示，尖銳物需進防刺穿容器。",
        ],
        "文件與現場查核回答": [
            "查核常看：SOP、教育訓練、稽核紀錄、缺失改善與追蹤是否能連到現場實作。",
            "同仁應能回答：自己的單位平常怎麼做、紀錄在哪裡、異常通報誰、改善後如何追蹤。",
            "回答時避免背誦式口號；用實際流程說明，並在不確定時回到院內政策與感染管制中心確認。",
        ],
    }
    return prompts.get(topic, [])


def main() -> None:
    topic_hits, used_files = collect()
    OUT.parent.mkdir(parents=True, exist_ok=True)

    lines: list[str] = [
        "# 115年感管查核臨床同仁重點",
        "",
        "本檔整理自 115 年感管查核相關資料，供 LINE 回答臨床同仁詢問查核、現場準備、常見感染管制作業觀念時使用。回答時請以臨床可執行的重點說明，不要輸出內部檔名、行政座位表、簽到或陪評安排。",
        "",
        "## 回答原則",
        "",
        "- 先回答同仁現場要怎麼做，再補充原因。",
        "- 遇到查核問題，不要只說「依規定」；應說出可執行行為，例如手部衛生、隔離標示、PPE、通報、採檢、清消、紀錄與通知流程。",
        "- 若問題涉及病人個資、感染診斷或抗藥菌狀態，提醒只做醫療照護必要範圍內揭露，避免公開標示疾病或菌名造成污名化。",
        "- 若牽涉院內最新政策、特殊感染症或查核當日指示，請提醒依院內政策規章與感染管制中心最新公告辦理。",
        "",
        "## 臨床同仁應知道的查核主題",
        "",
    ]

    for topic in TOPICS:
        lines.append(f"### {topic}")
        lines.append("")
        for item in manual_guidance(topic):
            lines.append(f"- {item}")
        for item in audit_prompts(topic):
            lines.append(f"- {item}")
        lines.append("")

    lines.extend(
        [
            "## 常見問法建議回答",
            "",
            "問：查核委員問手部衛生，我要怎麼回答？",
            "",
            "答：可以直接說明自己在照護前、無菌操作前、暴露體液風險後、接觸病人後、接觸病人周邊環境後都會執行手部衛生；戴手套前後也要洗手或乾洗手，手套不能取代手部衛生。",
            "",
            "問：查核時被問隔離病人要注意什麼？",
            "",
            "答：先確認是哪一種傳播途徑，再依接觸、飛沫、空氣或特殊防護準備病室、標示、PPE、器材專用或清消、檢查轉送通知與終期清潔。公開標示應以防護類型為主，不應揭露不必要的診斷或菌名。",
            "",
            "問：查核委員問抗藥菌病人怎麼照護？",
            "",
            "答：重點是接觸隔離、手部衛生、環境清消、共用器材清消或專用、檢查前通知接收單位、出院或轉床後終期清潔。若是既往陽性或再入院，依院內篩檢與解隔流程評估，不自行解除。",
            "",
            "問：同仁可不可以為特殊病人檢體、病歷、系統加醒目標記？",
            "",
            "答：標記應以病人安全與必要防護為目的，使用院內正式系統、檢驗單欄位或核准隔離標示。不要在公開區域或非必要文件寫 HIV、愛滋、VRE、CRE 等疾病或菌名；所有檢體本來就應依標準防護處理。",
            "",
            "問：查核委員問我不知道的細節怎麼辦？",
            "",
            "答：不要猜。可以回答自己會先依單位 SOP 與院內政策規章處理，並立即向單位主管、感染管制中心或相關專責單位確認；後續會補紀錄與改善措施。",
            "",
            "## LINE 問答關鍵字",
            "",
            "感管查核、115感管查核、查核委員、查核重點、查核準備、現場查核、自評表、佐證資料、手部衛生、隔離標示、PPE、N95、接觸隔離、飛沫隔離、空氣隔離、MDRO、VRE、CRE、抗生素管理、導尿管、中心導管、呼吸器、環境清潔、終期清潔、供應室、消毒滅菌、尖銳物、針扎、血液體液暴露、員工健康監測、法定傳染病通報、TOCC、HAI、THAS、病人安置、轉送檢查、感染性廢棄物、污衣。",
            "",
            "## 整理來源範圍",
            "",
            "本檔由 115 年感管查核資料夾內的查核作業手冊、查核基準及評分說明、自評表、全院宣導、護理部宣導、健康監測、供應室與查核結果相關文件彙整。LINE 回答時不要列出內部檔案路徑或檔名。",
            "",
        ]
    )

    OUT.write_text("\n".join(lines), encoding="utf-8")
    print(f"Wrote {OUT}")
    print(f"Used files: {len(used_files)}")
    for path in used_files[:30]:
        safe = str(path).encode("cp950", errors="ignore").decode("cp950", errors="ignore")
        print(safe)


if __name__ == "__main__":
    main()
